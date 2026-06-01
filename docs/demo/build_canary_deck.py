#!/usr/bin/env python3
"""Build CANary hackathon / demo presentation."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "CANary-BMS-Demo-Deck.pptx"

# Battery / industrial palette
BG_DARK = RGBColor(0x21, 0x29, 0x5C)
BG_LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
TEAL = RGBColor(0x02, 0x80, 0x90)
MINT = RGBColor(0x02, 0xC3, 0x9A)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x36, 0x45, 0x4F)
SLATE = RGBColor(0x5A, 0x6A, 0x7A)
CARD = RGBColor(0xE8, 0xED, 0xF2)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = CHARCOAL,
    align=PP_ALIGN.LEFT,
    font_name: str = "Calibri",
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font_name
    p.font.color.rgb = color


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    *,
    subtitle: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_LIGHT)

    # accent bar left
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    add_textbox(slide, Inches(0.55), Inches(0.45), Inches(9), Inches(0.7), title, size=34, bold=True, color=BG_DARK)

    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(1.05), Inches(9), Inches(0.5), subtitle, size=16, color=SLATE)

    top = Inches(1.55) if subtitle else Inches(1.25)
    box = slide.shapes.add_textbox(Inches(0.65), top, Inches(8.8), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(17)
        p.font.name = "Calibri"
        p.font.color.rgb = CHARCOAL
        p.space_after = Pt(14)


def add_card(slide, left, top, width, height, title: str, body: str, accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = CARD
    shape.line.width = Pt(1)

    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.2), top + Inches(0.22), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent
    dot.line.fill.background()

    add_textbox(slide, left + Inches(0.48), top + Inches(0.12), width - Inches(0.6), Inches(0.45), title, size=15, bold=True, color=BG_DARK)
    add_textbox(slide, left + Inches(0.22), top + Inches(0.55), width - Inches(0.44), height - Inches(0.65), body, size=13, color=SLATE)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s1, BG_DARK)
    glow = s1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(7.2), Inches(-0.8), Inches(3.5), Inches(3.5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = TEAL
    glow.fill.transparency = 0.75
    glow.line.fill.background()

    add_textbox(s1, Inches(0.75), Inches(1.8), Inches(8.5), Inches(0.5), "CANary", size=16, bold=True, color=MINT)
    add_textbox(s1, Inches(0.75), Inches(2.2), Inches(8.5), Inches(1.2), "BMS Validation Workbench", size=44, bold=True, color=WHITE)
    add_textbox(
        s1,
        Inches(0.75),
        Inches(3.55),
        Inches(7.5),
        Inches(1.2),
        "Agentic battery pack design with remote knowledge on Tigris\nand DevSecOps trust via Opsera",
        size=20,
        color=RGBColor(0xCA, 0xDC, 0xFC),
    )
    add_textbox(s1, Inches(0.75), Inches(6.35), Inches(8), Inches(0.4), "Hackathon Demo  ·  May 2026", size=14, color=SLATE)

    # --- Slide 2: Problem ---
    add_bullet_slide(
        prs,
        "The Problem",
        [
            "BMS design sits between safety and cost — wrong thresholds or wiring can mean recalls and weeks of rework.",
            "Knowledge lives in scattered PDFs (datasheets, UL/IEC standards) — not in the design itself.",
            "Schematics and protection rules drift apart; firmware constants don't match the block diagram.",
            "Every new pack (12s2p LFP, 16S NMC…) forces engineers to re-read the same docs from scratch.",
            "When validation asks “why is OVP 3.65 V?”, the answer is buried in a 200-page PDF — not traceable.",
        ],
        subtitle="Battery engineers don't lack tools — they lack connected, auditable design artifacts.",
    )

    # --- Slide 3: Pain cards ---
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s3, BG_LIGHT)
    add_textbox(s3, Inches(0.55), Inches(0.45), Inches(9), Inches(0.7), "Why engineers feel the pain", size=34, bold=True, color=BG_DARK)
    cards = [
        ("PDF-shaped knowledge", "Datasheets and standards aren't queryable at design time.", AMBER),
        ("Diagram ↔ rules gap", "The schematic says one thing; safety YAML says another.", RGBColor(0xE8, 0x5D, 0x4C)),
        ("No compounding memory", "Each project re-derives OVP, AFE pins, and topology from zero.", TEAL),
        ("Chatbots aren't enough", "LLM answers without schema, validation, or linked artifacts.", MINT),
    ]
    positions = [(0.55, 1.35), (5.05, 1.35), (0.55, 3.85), (5.05, 3.85)]
    for (left, top), (title, body, accent) in zip(positions, cards):
        add_card(s3, Inches(left), Inches(top), Inches(4.35), Inches(2.2), title, body, accent)

    # --- Slide 4: Solution ---
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s4, BG_DARK)
    add_textbox(s4, Inches(0.75), Inches(0.55), Inches(8.5), Inches(0.7), "Our Solution", size=36, bold=True, color=WHITE)
    add_textbox(
        s4,
        Inches(0.75),
        Inches(1.35),
        Inches(8.5),
        Inches(1.0),
        "CANary: describe the pack in natural language → agent produces validated artifacts → UI renders an interactive schematic.",
        size=20,
        color=RGBColor(0xCA, 0xDC, 0xFC),
    )
    steps = [
        ("1", "Engineer describes pack", "Topology, chemistry, parts, protection requirements"),
        ("2", "Agent reads Tigris wiki", "Datasheets, entities, protection concepts — cited, not guessed"),
        ("3", "Agent writes local files", "architecture.bms.json + safety_rules.yaml"),
        ("4", "Workbench renders SVG", "Pack drill-down, MCU rules inspector, schema validation"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.75 + i * 2.35)
        circle = s4.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, Inches(3.0), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = MINT
        circle.line.fill.background()
        add_textbox(s4, x, Inches(3.08), Inches(0.55), Inches(0.4), num, size=18, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_textbox(s4, x - Inches(0.15), Inches(3.75), Inches(2.1), Inches(0.5), title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s4, x - Inches(0.15), Inches(4.25), Inches(2.1), Inches(1.0), desc, size=12, color=SLATE, align=PP_ALIGN.CENTER)

    # --- Slide 5: Architecture split ---
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s5, BG_LIGHT)
    add_textbox(s5, Inches(0.55), Inches(0.45), Inches(9), Inches(0.7), "Architecture: local speed + remote brain", size=32, bold=True, color=BG_DARK)

    local_box = s5.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.35), Inches(4.2), Inches(5.2))
    local_box.fill.solid()
    local_box.fill.fore_color.rgb = WHITE
    local_box.line.color.rgb = TEAL
    local_box.line.width = Pt(2)
    add_textbox(s5, Inches(0.8), Inches(1.55), Inches(3.7), Inches(0.4), "Local workspace (fast UI)", size=18, bold=True, color=TEAL)
    add_textbox(
        s5,
        Inches(0.8),
        Inches(2.05),
        Inches(3.7),
        Inches(4.2),
        "• architecture.bms.json\n• safety_rules.yaml\n• React SVG renderer\n• JSON Schema validation\n• Agent SKILL + templates",
        size=15,
        color=CHARCOAL,
    )

    remote_box = s5.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.25), Inches(1.35), Inches(4.2), Inches(5.2))
    remote_box.fill.solid()
    remote_box.fill.fore_color.rgb = WHITE
    remote_box.line.color.rgb = MINT
    remote_box.line.width = Pt(2)
    add_textbox(s5, Inches(5.5), Inches(1.55), Inches(3.7), Inches(0.4), "Tigris (compounding knowledge)", size=18, bold=True, color=MINT)
    add_textbox(
        s5,
        Inches(5.5),
        Inches(2.05),
        Inches(3.7),
        Inches(4.2),
        "• raw/datasheets/*.pdf\n• wiki/index.md → entities, concepts\n• rtrvr.ai compiled sources\n• manifest.yaml catalog\n• Karpathy LLM Wiki pattern",
        size=15,
        color=CHARCOAL,
    )

    arrow = s5.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(4.55), Inches(3.6), Inches(0.55), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = AMBER
    arrow.line.fill.background()

    # --- Slide 6: Tigris ---
    add_bullet_slide(
        prs,
        "Tigris — BMS knowledge layer",
        [
            "Bucket: canary-bms-knowledge — 29+ wiki objects live today (entities, concepts, sources, PDFs).",
            "Agent always starts at wiki/index.md via Tigris MCP — users never say “read Tigris.”",
            "Immutable PDFs in raw/; agent-maintained markdown in wiki/ — knowledge compounds across sessions.",
            "rtrvr.ai batch ingest: vendor URLs → raw/rtrvr/*.json → compiled wiki/sources/*.md.",
            "Demo line: “Git holds the diagram; Tigris holds the brain.”",
        ],
        subtitle="Sponsor integration · Tigris Data",
    )

    # --- Slide 7: Opsera ---
    add_bullet_slide(
        prs,
        "Opsera — trust layer on the workbench",
        [
            "Separate from BMS design: security-scan, architecture-analyze, compliance-audit via MCP.",
            "Scans the CANary repo / workspace — not the battery pack — for critical/high findings.",
            "Reports saved to docs/opsera-scan/reports/ (security, architecture, SOC2 audit).",
            "Deep Agent + Cursor: make opsera-login once; scans run inside the workbench chat.",
            "Demo line: “Tigris validates the battery design; Opsera validates the software building it.”",
        ],
        subtitle="Sponsor integration · Opsera DevSecOps",
    )

    # --- Slide 8: Live demo ---
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s8, BG_DARK)
    add_textbox(s8, Inches(0.75), Inches(0.55), Inches(8.5), Inches(0.7), "Live Demo Flow", size=36, bold=True, color=WHITE)
    demo_steps = [
        "Show Tigris bucket — index, BQ76952 PDF, rtrvr source pages",
        "Prompt: “Design 12s2p LFP with BQ76952 + STM32F407 — use remote wiki for thresholds”",
        "Watch agent: tigris_get_object → write architecture + safety rules",
        "Diagram tab: pack view → drill into BMS board → inspect MCU protection rules",
        "Follow-up: tighten thermal limits — rules update with wiki citation",
    ]
    for i, step in enumerate(demo_steps):
        y = Inches(1.45 + i * 1.05)
        badge = s8.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), y, Inches(0.45), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        add_textbox(s8, Inches(0.75), y + Inches(0.06), Inches(0.45), Inches(0.35), str(i + 1), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s8, Inches(1.4), y + Inches(0.05), Inches(8), Inches(0.6), step, size=17, color=RGBColor(0xCA, 0xDC, 0xFC))

    # --- Slide 9: Results ---
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s9, BG_LIGHT)
    add_textbox(s9, Inches(0.55), Inches(0.45), Inches(9), Inches(0.7), "What we built", size=34, bold=True, color=BG_DARK)
    stats = [
        ("29+", "Wiki pages\non Tigris"),
        ("4.1 MB", "Real TI\nBQ76952 PDF"),
        ("2", "Core artifacts\n(JSON + YAML)"),
        ("10", "Tigris MCP\ntools wired"),
    ]
    for i, (num, label) in enumerate(stats):
        x = Inches(0.55 + i * 2.35)
        card = s9.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(1.4), Inches(2.1), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = CARD
        add_textbox(s9, x, Inches(1.65), Inches(2.1), Inches(0.8), num, size=36, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_textbox(s9, x, Inches(2.5), Inches(2.1), Inches(0.8), label, size=13, color=SLATE, align=PP_ALIGN.CENTER)

    add_textbox(
        s9,
        Inches(0.65),
        Inches(3.75),
        Inches(8.8),
        Inches(2.8),
        "Shipped: FastAPI Deep Agent workbench · React SVG schematic renderer · Schema-backed BMS validation\n"
        "· Tigris remote wiki + MCP · rtrvr.ai ingest pipeline · Opsera security & compliance reports\n"
        "· Natural-language pack design in one turn with cited thresholds (not hardcoded in UI)",
        size=15,
        color=CHARCOAL,
    )

    # --- Slide 10: Roadmap ---
    add_bullet_slide(
        prs,
        "What's next",
        [
            "Pack-level simulation — fault injection, CAN traces, scenario runners.",
            "Validation reports — requirements traceability and verification outputs.",
            "More datasheets ingested via rtrvr + Tigris (UL 2580, IEC 62619, chemistry guides).",
            "Formal certification workflow — wiki-linked evidence for audit questions.",
        ],
        subtitle="Hackathon proves design + knowledge; production adds simulation and sign-off.",
    )

    # --- Slide 11: Close ---
    s11 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s11, BG_DARK)
    add_textbox(s11, Inches(0.75), Inches(2.2), Inches(8.5), Inches(1.0), "Design safe battery systems\nthat learn from every project.", size=40, bold=True, color=WHITE)
    add_textbox(
        s11,
        Inches(0.75),
        Inches(4.0),
        Inches(8.5),
        Inches(0.8),
        "CANary · Tigris remote wiki · Opsera DevSecOps · Retriever AI ingest",
        size=18,
        color=MINT,
    )
    add_textbox(s11, Inches(0.75), Inches(5.8), Inches(8.5), Inches(0.5), "Questions?", size=28, bold=True, color=AMBER)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
